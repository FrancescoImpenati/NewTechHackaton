"""Faithful-ish raster preview of a .pptx by reading real shape geometry."""
import io, textwrap, sys
from pathlib import Path
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle
from PIL import Image
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

EMU=914400.0
SW,SH=13.333,7.5
def IN(e): return e/EMU if e is not None else 0

def rgb(color):
    try: return "#%02x%02x%02x"%(color[0],color[1],color[2])
    except Exception: return None

def shape_fill(sh):
    try:
        if sh.fill.type is not None and sh.fill.type==1:  # solid
            return rgb(sh.fill.fore_color.rgb)
    except Exception: pass
    return None
def shape_line(sh):
    try:
        c=sh.line.color.rgb; return rgb(c)
    except Exception: return None

def render(path, outdir):
    prs=Presentation(path); outdir=Path(outdir); outdir.mkdir(exist_ok=True)
    files=[]
    for si,slide in enumerate(prs.slides,1):
        fig=plt.figure(figsize=(SW,SH),dpi=100); ax=fig.add_axes([0,0,1,1])
        ax.set_xlim(0,SW); ax.set_ylim(0,SH); ax.invert_yaxis(); ax.axis("off")
        ax.add_patch(Rectangle((0,0),SW,SH,facecolor="#ffffff",zorder=0))
        for zi,sh in enumerate(slide.shapes):
            Z=10+zi*3
            x,y,w,h=IN(sh.left),IN(sh.top),IN(sh.width),IN(sh.height)
            try: st=sh.shape_type
            except Exception: st=None
            if st==MSO_SHAPE_TYPE.PICTURE:
                try:
                    img=Image.open(io.BytesIO(sh.image.blob))
                    ax.imshow(img,extent=(x,x+w,y+h,y),zorder=Z,aspect="auto")
                except Exception: pass
                continue
            fc=shape_fill(sh); lc=shape_line(sh)
            if fc or lc:
                rounded = "ROUND" in str(getattr(sh,'auto_shape_type',''))
                if rounded:
                    p=FancyBboxPatch((x+0.03,y+0.03),max(w-0.06,0.01),max(h-0.06,0.01),
                        boxstyle="round,pad=0.03,rounding_size=0.08",
                        facecolor=fc or "none",edgecolor=lc or "none",
                        linewidth=(sh.line.width/12700 if lc else 0),zorder=Z+1)
                else:
                    p=Rectangle((x,y),w,h,facecolor=fc or "none",edgecolor=lc or "none",
                        linewidth=(sh.line.width/12700 if lc else 0),zorder=Z+1)
                ax.add_patch(p)
            # text
            if sh.has_text_frame and sh.text_frame.text.strip():
                tf=sh.text_frame
                # gather first run formatting per paragraph
                yy=y+0.06; anchor=str(tf.vertical_anchor)
                paras=[p for p in tf.paragraphs]
                # estimate total height for middle/bottom anchor
                def psize(p):
                    for r in p.runs:
                        if r.font.size: return r.font.size.pt
                    return 14
                total=sum((psize(p)/72.0*1.25)+0.02 for p in paras if p.text.strip())
                if "MIDDLE" in anchor: yy=y+(h-total)/2+psize(paras[0])/72.0*0.8
                for p in paras:
                    if not p.runs:
                        if p.text.strip(): pass
                        else: continue
                    sz=psize(p); txt="".join(r.text for r in p.runs)
                    if not txt.strip(): continue
                    col="#222222"; bold=False
                    for r in p.runs:
                        if r.text.strip():
                            col=rgb(r.font.color.rgb) if (r.font.color and r.font.color.type) else col
                            bold=bool(r.font.bold); break
                    al=p.alignment
                    ha="left"; tx=x+0.08
                    if al==PP_ALIGN.CENTER: ha="center"; tx=x+w/2
                    elif al==PP_ALIGN.RIGHT: ha="right"; tx=x+w-0.08
                    # wrap
                    charw=sz*0.0085  # inches per char approx for DejaVu
                    maxch=max(int((w-0.16)/max(charw,0.001)),4)
                    lines=[]
                    for seg in txt.split("\n"):
                        lines+=textwrap.wrap(seg,maxch) or [""]
                    for ln in lines:
                        ax.text(tx,yy+sz/72.0*0.8,ln,fontsize=sz,color=col or "#222",
                            ha=ha,va="baseline",weight="bold" if bold else "normal",
                            family="DejaVu Sans",zorder=Z+2)
                        yy+=sz/72.0*1.22
        f=outdir/f"slide-{si:02d}.png"; fig.savefig(f,dpi=100); plt.close(fig); files.append(str(f))
    print("\n".join(files))

if __name__=="__main__":
    render(sys.argv[1], sys.argv[2] if len(sys.argv)>2 else "qa")
