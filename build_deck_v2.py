"""Build EWS_pitch_deck_v2.pptx — 15 slides, 16:9, trading-terminal design."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent
A = ROOT / "outputs" / "deck_assets"
OUT = ROOT / "outputs" / "EWS_pitch_deck_v2.pptx"

# palette
DARK=RGBColor(0x0B,0x12,0x20); PANEL=RGBColor(0x16,0x22,0x3B); PANEL2=RGBColor(0x22,0x32,0x4F)
INK=RGBColor(0xEA,0xF0,0xF8); MUTE=RGBColor(0x93,0xA6,0xC2); ICE=RGBColor(0x5E,0x84,0xB5)
ACCENT=RGBColor(0xF2,0xB1,0x34); GREEN=RGBColor(0x3F,0xB3,0x6B); RED=RGBColor(0xE2,0x50,0x3B)
BLUE=RGBColor(0x4E,0x8F,0xD0); GREY=RGBColor(0x8A,0x93,0xA6)
LIGHT=RGBColor(0xF3,0xF6,0xFA); CARD=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xD7,0xDF,0xEA)
NAVY=RGBColor(0x14,0x21,0x3D); DMUTE=RGBColor(0x5B,0x6B,0x86)
FONT="DejaVu Sans"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
def slide(): return prs.slides.add_slide(prs.slide_layouts[6])

def bg(s,color):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=color
    r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return r

def rrect(s,x,y,w,h,fill,line=None,lw=0.0,rad=0.06):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is not None: sh.line.color.rgb=line; sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    sh.shadow.inherit=False
    try: sh.adjustments[0]=rad
    except Exception: pass
    return sh

def dot(s,x,y,d,color):
    o=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb=color; o.line.fill.background(); o.shadow.inherit=False; return o

def txt(s,string,x,y,w,h,size,color,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,
        italic=False,font=FONT,spacing=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=align; p.line_spacing=spacing
    r=p.add_run(); r.text=string; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font; f.color.rgb=color
    return tb

def kicker(s,string,x=0.62,y=0.5,color=ICE):
    txt(s," ".join(string.upper()),x,y,11.0,0.3,11.5,color,bold=True)

def title(s,string,x=0.62,y=0.86,color=NAVY,size=30,w=11.5):
    dot(s,x,y+0.10,0.13,ACCENT)
    txt(s,string,x+0.28,y,w,0.7,size,color,bold=True)

def bullets(s,items,x,y,w,size=15,tc=NAVY,mk=ACCENT,gap=10,sp=1.05,h=4.6):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.line_spacing=sp; p.space_after=Pt(gap)
        r1=p.add_run(); r1.text="●  "; r1.font.size=Pt(size-3); r1.font.color.rgb=mk; r1.font.name=FONT
        r2=p.add_run(); r2.text=it; r2.font.size=Pt(size); r2.font.color.rgb=tc; r2.font.name=FONT
    return tb

def stat(s,num,label,x,y,w,nc=ACCENT,lc=MUTE,nsize=44,lsize=11.5,align=PP_ALIGN.LEFT):
    txt(s,num,x,y,w,0.8,nsize,nc,bold=True,align=align)
    txt(s,label,x,y+nsize/64.0+0.02,w,0.5,lsize,lc,align=align)

def pagenum(s,n,color=DMUTE):
    txt(s,f"{n:02d} / 15",11.7,7.05,1.4,0.3,10,color,align=PP_ALIGN.RIGHT)

def pic(s,name,x,y,w=None,h=None):
    kw={}
    if w: kw["width"]=Inches(w)
    if h: kw["height"]=Inches(h)
    return s.shapes.add_picture(str(A/name),Inches(x),Inches(y),**kw)

# ============================== S1 TITLE ==============================
s=slide(); bg(s,DARK)
# decorative node cluster top-right (network motif)
import random; random.seed(7)
pts=[(11.6,0.7),(12.3,1.1),(11.9,1.7),(12.6,2.0),(11.3,1.4),(12.1,2.6)]
for i in range(len(pts)-1):
    ln=s.shapes.add_connector(2,Inches(pts[i][0]+0.07),Inches(pts[i][1]+0.07),Inches(pts[i+1][0]+0.07),Inches(pts[i+1][1]+0.07))
    ln.line.color.rgb=PANEL2; ln.line.width=Pt(1.2)
for i,(px,py) in enumerate(pts):
    dot(s,px,py,0.16 if i%2 else 0.12, ACCENT if i==0 else (ICE if i%2 else BLUE))
kicker(s,"ML for Fintech  ·  Politecnico di Milano — GSoM",x=0.9,y=1.5,color=ICE)
txt(s,"Early Warning System",0.9,2.25,11,1.0,46,INK,bold=True)
txt(s,"for Risk-Off Detection",0.9,3.1,11,1.0,46,INK,bold=True)
txt(s,"A Quant Strategy: 1.5× Equity + Macro-Aware Safe-Haven Routing",0.92,4.25,11,0.6,19,ACCENT,bold=True)
txt(s,"Ercelli  ·  Galli  ·  Impenati  ·  Mazzini",0.92,5.5,11,0.4,15,INK)
txt(s,"Bloomberg weekly  ·  2000–2021  ·  n ≈ 1,107  ·  2026",0.92,5.95,11,0.4,12.5,MUTE)
dot(s,0.62,2.45,0.14,ACCENT)

# ============================== S2 PROBLEM ==============================
s=slide(); bg(s,DARK)
kicker(s,"The problem & business case",color=ICE)
title(s,"Avoiding the worst weeks beats picking the best",color=INK)
bullets(s,[
 "Crashes cluster — buy-and-hold pays for every one in full.",
 "Skipping a handful of weeks compounds into far higher Sharpe & Calmar.",
 "Our answer is a Quant Strategy: aggressive in calm, defensive in stress.",
 "Nowcasting, not forecasting — recognise stress at birth, don't predict it.",
],0.62,2.05,6.1,size=15.5,tc=INK,mk=ACCENT,gap=14)
rrect(s,7.05,1.95,5.7,4.45,PANEL,rad=0.04)
pic(s,"s2_drawdown.png",7.3,2.7,w=5.2,h=3.5)
stat(s,"−27%","MSCI buy-&-hold COVID drawdown",7.35,2.05,4,nc=RED,lc=MUTE,nsize=26,lsize=11)
pagenum(s,2)

# ============================== S3 DATA ==============================
s=slide(); bg(s,LIGHT)
kicker(s,"Data & target",color=ICE)
title(s,"Weekly Bloomberg, two decades, one binary signal")
for i,(n,l,c) in enumerate([("1,111","raw weekly obs · 2000–2021",NAVY),
                            ("43 → 56","raw → engineered features",NAVY),
                            ("21.3%","risk-off prevalence (Y=1)",RED)]):
    x=0.62+i*2.55
    rrect(s,x,1.95,2.35,1.25,CARD,LINE,1.0,rad=0.10)
    stat(s,n,l,x+0.22,2.1,2.0,nc=ACCENT if i<2 else RED,lc=DMUTE,nsize=30,lsize=10.5)
bullets(s,[
 "Asset classes: equity · FX · gold · oil · sovereign yields · credit · MBS.",
 "Risk-off clusters: Dot-com 2000–02 · GFC 2008 · Euro 2011 · COVID 2020.",
],0.62,3.45,7.4,size=14.5,tc=NAVY,mk=ICE,gap=8,h=1.3)
rrect(s,8.25,1.95,4.5,1.25,RGBColor(0xFD,0xF1,0xE0),ACCENT,1.0,rad=0.10)
txt(s,"The CV challenge",8.5,2.08,4.0,0.3,12,RED,bold=True)
txt(s,"2013 · 2017 · 2019 carry ZERO positives — folds must be weighted by event count.",
    8.5,2.42,4.05,0.7,12,NAVY,spacing=1.05)
rrect(s,0.62,4.55,12.1,2.4,CARD,LINE,1.0,rad=0.04)
pic(s,"s3_clusters.png",0.95,4.8,w=11.4,h=1.95)
pagenum(s,3)

# ============================== S4 FEATURES ==============================
s=slide(); bg(s,LIGHT)
kicker(s,"Feature engineering")
title(s,"From 43 raw series to 56 stationary, decorrelated predictors")
cards=[("Stationarity by family","Log-returns for prices · first differences for yields · level passthrough for VIX & surprise indices.",ICE),
       ("8 macro spreads","Term, sovereign and credit (log-ratio) spreads — each as level and 4-week change.",BLUE),
       ("Collinearity cleanup","Drop one of every pair with |corr| > 0.9 → a lean, non-redundant 56-feature set.",GREY)]
for i,(t,d,c) in enumerate(cards):
    x=0.62+i*4.05
    rrect(s,x,2.05,3.8,2.35,CARD,LINE,1.0,rad=0.06)
    dot(s,x+0.28,2.32,0.16,c); txt(s,t,x+0.6,2.26,3.0,0.4,14.5,NAVY,bold=True)
    txt(s,d,x+0.3,2.85,3.25,1.4,12,DMUTE,spacing=1.1)
rrect(s,0.62,4.75,12.1,1.9,NAVY,rad=0.05)
dot(s,1.0,5.15,0.2,ACCENT)
txt(s,"Hero feature — equity_bond_corr_13w",1.4,5.0,7,0.4,16,ACCENT,bold=True)
txt(s,"The 13-week stock–bond return correlation separates liquidity-driven crises (correlation turns positive, diversification fails) from inflation-driven regimes — a discriminating axis the price and yield features alone miss.",
    1.4,5.5,10.9,1.1,13.5,INK,spacing=1.12)
pagenum(s,4)

# ============================== S5 ERRORS ==============================
s=slide(); bg(s,LIGHT)
kicker(s,"Methodological rigour",color=RED)
title(s,'“You must not fool yourself” — the traps we engineered against')
grid=[("Scaler fit on the full dataset","Fit StandardScaler on each fold's TRAIN only"),
      ("Random split on a time series","Purged expanding walk-forward + 4-week embargo"),
      ("Threshold tuned on the test set","Tuned on CV folds only; 2019–21 sealed holdout"),
      ("Accuracy on imbalanced labels","F1 & AUC-PR on the anomaly class, weighted by n_pos")]
for i,(trap,fix) in enumerate(grid):
    cx=0.62+(i%2)*6.15; cy=2.1+(i//2)*2.35
    rrect(s,cx,cy,5.9,2.1,CARD,LINE,1.0,rad=0.06)
    dot(s,cx+0.28,cy+0.32,0.16,RED); txt(s,"TRAP",cx+0.56,cy+0.24,2,0.3,11,RED,bold=True)
    txt(s,trap,cx+0.3,cy+0.6,5.3,0.5,14,NAVY,bold=True,spacing=1.0)
    txt(s,"↓",cx+0.3,cy+1.05,5.3,0.3,13,DMUTE)
    dot(s,cx+0.28,cy+1.5,0.16,GREEN); txt(s,"OUR FIX",cx+0.56,cy+1.42,2,0.3,11,GREEN,bold=True)
    txt(s,fix,cx+0.3,cy+1.72,5.4,0.4,12.5,NAVY,spacing=1.0)
pagenum(s,5)

# ============================== S6 METHODOLOGY ==============================
s=slide(); bg(s,LIGHT)
kicker(s,"Methodology & economic rationale")
title(s,"Lead with why, not just what")
flow=["Stationarity","Spreads","Walk-forward","4 detectors","Ensemble","Routing","Backtest"]
fx=0.62; fw=12.1; step=fw/len(flow)
for i,st in enumerate(flow):
    x=fx+i*step
    rrect(s,x,2.0,step-0.18,0.62,PANEL if i%2 else PANEL2,rad=0.18)
    txt(s,st,x,2.13,step-0.18,0.4,10.5,INK,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if i<len(flow)-1: txt(s,"›",x+step-0.16,2.05,0.2,0.5,16,ACCENT,bold=True)
rows=[("Stationarity","prices follow a random walk; their returns are stationary — ADF-clean inputs."),
      ("Purged walk-forward + 4-week embargo","the embargo gap kills 4-week-lookback leakage between train and validation."),
      ("Per-fold StandardScaler on train only","validation statistics never bleed into training — no cross-fold leakage."),
      ("Ledoit-Wolf shrinkage","56 features vs ~254 normals → the sample covariance is ill-conditioned; shrinkage tames noisy eigenvectors.")]
y=3.15
for t,d in rows:
    rrect(s,0.62,y,12.1,0.78,CARD,LINE,0.75,rad=0.05)
    dot(s,0.92,y+0.31,0.17,ACCENT)
    tb=s.shapes.add_textbox(Inches(1.25),Inches(y+0.13),Inches(11.2),Inches(0.55))
    tf=tb.text_frame; tf.word_wrap=True; tf.margin_top=0; tf.margin_left=0
    p=tf.paragraphs[0]; p.line_spacing=1.05
    r1=p.add_run(); r1.text=t+" — "; r1.font.bold=True; r1.font.size=Pt(13.5); r1.font.color.rgb=NAVY; r1.font.name=FONT
    r2=p.add_run(); r2.text=d; r2.font.size=Pt(13.5); r2.font.color.rgb=DMUTE; r2.font.name=FONT
    y+=0.9
pagenum(s,6)

# ============================== S7 BACKTEST ==============================
s=slide(); bg(s,LIGHT)
kicker(s,"The honest backtest")
title(s,"Five expanding folds, weighted by what's rare")
hdr=["Fold","Validation window","Positives","Crisis"]
colx=[0.62,1.7,5.1,6.4]; colw=[1.0,3.3,1.2,2.0]
rrect(s,0.62,2.05,8.1,0.5,NAVY,rad=0.04)
for c,(hx,hh) in enumerate(zip(colx,hdr)):
    txt(s,hh,hx+0.1,2.13,colw[c],0.35,12,INK,bold=True)
fr=[("1","2007-01 → 2009-12","53","GFC 2008",True),
    ("2","2010-02 → 2012-12","64","Euro 2011",True),
    ("3","2013-01 → 2014-12","2","Taper 2013",False),
    ("4","2015-02 → 2016-12","10","China–Oil 2015-16",False),
    ("5","2017-01 → 2018-12","6","Q4-2018 selloff",False)]
y=2.6
for f,win,pos,cr,prim in fr:
    rrect(s,0.62,y,8.1,0.52,CARD if prim else RGBColor(0xEE,0xF1,0xF6),LINE,0.75,rad=0.04)
    txt(s,f,0.72,y+0.1,1,0.3,12.5,NAVY,bold=True)
    txt(s,win,1.8,y+0.1,3.3,0.3,12,DMUTE)
    txt(s,pos,5.2,y+0.1,1.2,0.3,12.5,ACCENT if prim else RED,bold=True)
    txt(s,cr,6.5,y+0.1,2.1,0.3,12,NAVY)
    y+=0.6
bullets(s,[
 "Folds 1–2 hold 117 of 135 positives — they carry the signal.",
 "CV metrics aggregated weighted by n_pos; thin folds 3–5 down-weighted.",
 "Sealed 2019–2021 test holdout: 121 weeks, includes COVID.",
],9.0,2.6,3.8,size=13.5,tc=NAVY,mk=ICE,gap=14,h=3)
rrect(s,9.0,4.45,3.75,2.2,NAVY,rad=0.05)
stat(s,"87%","of all positives sit in folds 1–2",9.3,4.75,3.2,nc=ACCENT,lc=MUTE,nsize=46,lsize=12.5)
txt(s,"117 of 135 risk-off weeks — why we weight folds by n_pos.",9.3,5.9,3.25,0.6,11.5,INK,spacing=1.1)
pagenum(s,7)

# ============================== S8 BRIDGE (dark) ==============================
s=slide(); bg(s,DARK)
kicker(s,"Macro intuition",color=ACCENT)
txt(s,"Knowing THAT there's a crisis isn't enough —",0.9,1.7,11.5,0.7,30,INK,bold=True)
txt(s,"you need to know WHERE to hide.",0.9,2.5,11.5,0.7,30,ACCENT,bold=True)
txt(s,"The right haven depends on the nature of the shock. We score three of them.",
    0.92,3.5,11,0.5,15,MUTE)
doms=[("USD","Funding stress / dollar squeeze","LIBOR-OIS, DXY, VRP, 10Y move, US-vs-world",BLUE),
      ("GOLD","Real-yield · inflation · fiat distrust","real yields, DXY, JPY, stock-bond corr, gold/oil",ACCENT),
      ("MBS","Moderate-stress carry rule","active in mild stress, blocked in acute crises",GREY)]
for i,(t,sub,d,c) in enumerate(doms):
    x=0.9+i*3.95
    rrect(s,x,4.3,3.7,2.4,PANEL,rad=0.06)
    dot(s,x+0.3,4.62,0.22,c); txt(s,t,x+0.66,4.55,3,0.4,18,c,bold=True)
    txt(s,sub,x+0.3,5.15,3.2,0.5,12.5,INK,bold=True,spacing=1.0)
    txt(s,d,x+0.3,5.75,3.25,0.8,11,MUTE,spacing=1.1)
pagenum(s,8,color=DMUTE)

# ============================== S9 HERO NETWORK ==============================
s=slide(); bg(s,LIGHT)
kicker(s,"The EWS trigger map")
title(s,"14 signals → 3 domains → router → 4 regimes")
pic(s,"s9_network.png",0.5,1.58,w=12.4,h=5.12)
txt(s,"dxy_chg4w is shared by USD (+) and Gold (−) — one variable, opposite economic readings.",
    0.7,6.86,12,0.3,11,DMUTE,italic=True)
pagenum(s,9)

# ============================== S10 DECISION MATRIX ==============================
s=slide(); bg(s,LIGHT)
kicker(s,"Decision matrix")
title(s,"From signals to allocation")
rules=[("Ensemble = risk-on","→","LEVERED_EQUITY  ·  1.5× equity, −0.5× cash",GREEN),
       ("risk-off  ·  USD sub-score > 1.5  ·  dxy_chg4w > 0","→","CASH_USD",BLUE),
       ("risk-off  ·  Gold sub-score > 0.5","→","GOLD",ACCENT),
       ("risk-off  ·  MBS rule active","→","MBS",GREY),
       ("otherwise","→","CASH_USD  (safe default)",BLUE)]
y=2.15
for cond,arr,act,c in rules:
    rrect(s,0.62,y,8.5,0.74,CARD,LINE,1.0,rad=0.08)
    txt(s,cond,0.85,y+0.2,6.0,0.4,13,NAVY,anchor=MSO_ANCHOR.MIDDLE)
    rrect(s,6.9,y+0.13,2.0,0.48,c,rad=0.2)
    txt(s,act.split("  ·")[0] if "·" in act else act,6.9,y+0.2,2.0,0.34,11,DARK if c in (ACCENT,GREEN,GREY) else INK,bold=True,align=PP_ALIGN.CENTER)
    y+=0.86
rrect(s,9.4,2.15,3.3,4.0,NAVY,rad=0.05)
txt(s,"Thresholds",9.65,2.4,3,0.4,14,ACCENT,bold=True)
stat(s,"1.5","τ  USD",9.65,2.95,3,nc=INK,lc=MUTE,nsize=40,lsize=12)
stat(s,"0.5","τ  Gold",9.65,4.15,3,nc=INK,lc=MUTE,nsize=40,lsize=12)
txt(s,"Tuned via walk-forward Calmar grid (duration-weighted median).",
    9.65,5.35,2.85,0.9,11.5,MUTE,spacing=1.12)
pagenum(s,10)

# ============================== S11 MODELS ==============================
s=slide(); bg(s,LIGHT)
kicker(s,"The four detectors")
title(s,"Four inductive biases, one interface")
ms=[("MVG · Ledoit-Wolf","parametric — interpretable Mahalanobis distance",ICE),
    ("One-Class SVM (RBF)","flexible non-linear boundary; grid ν × γ",BLUE),
    ("Autoencoder 56-…-6-…-56","non-linear reconstruction error; dropout 0.15",GREEN),
    ("Isolation Forest","200 trees; robust, scale-invariant",GREY)]
y=2.15
for t,d,c in ms:
    rrect(s,0.62,y,6.0,0.92,CARD,LINE,1.0,rad=0.07)
    dot(s,0.9,y+0.36,0.18,c); txt(s,t,1.22,y+0.16,5.2,0.4,14,NAVY,bold=True)
    txt(s,d,1.22,y+0.52,5.2,0.35,11.5,DMUTE)
    y+=1.05
txt(s,"All fit on Y = 0 normals only.",0.62,6.55,6,0.3,12,DMUTE,italic=True)
rrect(s,6.95,2.05,5.8,3.55,PANEL,rad=0.04)
pic(s,"s11_models.png",7.15,2.2,w=5.4,h=3.25)
rrect(s,6.95,5.8,5.8,0.95,RGBColor(0xFD,0xF1,0xE0),ACCENT,1.0,rad=0.08)
txt(s,"AE wins AUC-PR (0.78); MVG best F1 — diversity that the ensemble exploits.",
    7.2,5.95,5.4,0.7,12.5,NAVY,bold=True,spacing=1.1)
pagenum(s,11)

# ============================== S12 ENSEMBLE ==============================
s=slide(); bg(s,LIGHT)
kicker(s,"Ensemble & error correlation")
title(s,"A real, but honest, lift")
bullets(s,[
 "Raw scores → empirical percentiles vs train normals (comparable scale).",
 "Three modes: hard vote (≥3/4), soft mean, soft median.",
 "Winner — soft median: test F1 = 0.647, beating the best single model (MVG, 0.635).",
 "Errors are correlated but far from identical → a genuine yet modest gain.",
 "To push further: diversify INPUTS — rates-only / equity-only / FX-only models.",
],0.62,2.1,6.2,size=14,tc=NAVY,mk=ACCENT,gap=11,h=4.2)
rrect(s,7.05,2.05,5.7,3.25,PANEL,rad=0.04)
pic(s,"s12_errcorr.png",7.85,2.25,w=4.1,h=2.85)
rrect(s,7.05,5.5,5.7,1.25,NAVY,rad=0.06)
stat(s,"0.507","mean pairwise error correlation (folds 1–2)",7.35,5.62,5.2,nc=ACCENT,lc=MUTE,nsize=34,lsize=11)
pagenum(s,12)

# ============================== S13 RESULTS ==============================
s=slide(); bg(s,LIGHT)
kicker(s,"Results — backtest")
title(s,"Growth of $1 on the sealed test holdout")
rrect(s,0.62,1.95,8.7,4.9,CARD,LINE,1.0,rad=0.04)
pic(s,"s13_equity.png",0.85,2.2,w=8.25,h=4.42)
for i,(n,l,c) in enumerate([("1.87×","EWS final — highest of 6",ACCENT),
                            ("22 wks","COVID recovery (MSCI: 35)",GREEN),
                            ("GOLD","haven during the crash",ACCENT)]):
    y=2.0+i*1.62
    rrect(s,9.55,y,3.2,1.42,NAVY,rad=0.06)
    stat(s,n,l,9.8,y+0.22,2.8,nc=c,lc=MUTE,nsize=34,lsize=11)
pagenum(s,13)

# ============================== S14 METRICS ==============================
s=slide(); bg(s,LIGHT)
kicker(s,"Metrics & the drawdown nuance")
title(s,"Best risk-adjusted return of all six")
mets=[("1.45","Sharpe"),("2.22","Sortino"),("1.39","Calmar"),("0.31","CAGR")]
for i,(n,l) in enumerate(mets):
    x=0.62+i*3.05
    rrect(s,x,2.0,2.85,1.5,NAVY,rad=0.06)
    stat(s,n,l,x+0.3,2.2,2.4,nc=ACCENT,lc=MUTE,nsize=44,lsize=13)
txt(s,"EWS wins all four versus every benchmark.",0.62,3.65,7,0.3,12.5,GREEN,bold=True)
rrect(s,0.62,4.15,6.0,2.55,PANEL,rad=0.04)
pic(s,"s14_dd.png",0.85,4.35,w=5.55,h=2.1)
rrect(s,6.95,4.15,5.8,2.55,RGBColor(0xFD,0xF1,0xE0),ACCENT,1.0,rad=0.05)
txt(s,"The honest nuance",7.2,4.35,5,0.4,14,RED,bold=True)
bullets(s,[
 "EWS max DD −22% vs MSCI −28%; COVID crash −21% vs −27%.",
 "Always-defensive blends (Butterfly / Permanent / Risk Parity) show smaller ABSOLUTE drawdowns — they are never fully invested.",
 "EWS is a Quant Strategy (best risk-adjusted return), not a risk-management overlay.",
],7.2,4.8,5.35,size=12.5,tc=NAVY,mk=RED,gap=8,h=1.9)
pagenum(s,14)

# ============================== S15 LIMITATIONS + ROADMAP (dark) ==============================
s=slide(); bg(s,DARK)
kicker(s,"Limitations & extensions",color=ACCENT)
title(s,"What we'd test next — tied to what's missing",color=INK)
txt(s,"Limitations",0.7,2.0,4,0.4,15,RED,bold=True)
lims=["Sparse labels post-2013 (thin folds).","Training is liquidity/credit-driven — no 2022 inflation shock.",
      "MSCI World = equal-weight DM proxy.","Credit = log-ratio of TR indices, not OAS.",
      "Static t-costs; 1-week execution lag."]
bullets(s,lims,0.7,2.5,5.4,size=12.5,tc=INK,mk=RED,gap=8,h=4)
txt(s,"Roadmap",6.6,2.0,4,0.4,15,GREEN,bold=True)
road=[("Signal diversity","ensembles over DISJOINT feature sets to break the 0.507 error correlation"),
      ("Regime modeling","HMM with 2–3 states for explicit switching"),
      ("Frequency","daily signals — the 1-week lag cost us the first COVID week"),
      ("Asset universe","REITs, EM bonds, commodity baskets"),
      ("OOS stress","2022 inflation, 2023-24 geopolitical, 2025 tariff regime")]
y=2.5
for t,d in road:
    dot(s,6.6,y+0.05,0.15,GREEN); txt(s,t,6.86,y-0.03,5.8,0.35,12.5,INK,bold=True)
    txt(s,d,6.86,y+0.27,5.9,0.45,11,MUTE,spacing=1.05)
    y+=0.82
pagenum(s,15,color=DMUTE)

prs.save(OUT); print("saved", OUT, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
print("slides:", len(prs.slides._sldIdLst))
